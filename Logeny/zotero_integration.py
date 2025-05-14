import os
import re
import shutil
import sqlite3
import uuid
from pathlib import Path
from datetime import datetime


import numpy as np
import pandas as pd
from sentence_transformers import SentenceTransformer

from extract_text import read_text_file  # We'll use your existing read_text_file() here.

###############################################################################
# 1) ZOTERO CORE DB LOGIC
###############################################################################

def get_all_collections(db_path):
    """
    Returns a sorted list of all collection names from the Zotero database.
    """
    if not os.path.exists(db_path):
        raise FileNotFoundError(f"Zotero database not found: {db_path}")

    conn = sqlite3.connect(db_path)
    c = conn.cursor()
    c.execute("SELECT collectionName FROM collections")
    rows = c.fetchall()
    conn.close()

    collection_names = [r[0] for r in rows if r[0] is not None]
    return sorted(set(collection_names))

def get_collection_items_metadata(db_path, collection_name, require_attachment=False):
    """
    Retrieves metadata (title, authors, year, key) for items in the given collection.
    If require_attachment=True, skip items that have no 'key' (file path).
    """
    if not os.path.exists(db_path):
        raise FileNotFoundError(f"Zotero database not found: {db_path}")

    conn = sqlite3.connect(db_path)
    c = conn.cursor()

    # 1) Find the collectionID
    c.execute("SELECT collectionID FROM collections WHERE collectionName=?", (collection_name,))
    row = c.fetchone()
    if not row:
        conn.close()
        return []
    collection_id = row[0]

    # 2) itemIDs in that collection
    c.execute("SELECT itemID FROM collectionItems WHERE collectionID=?", (collection_id,))
    item_ids = [r[0] for r in c.fetchall()]
    if not item_ids:
        conn.close()
        return []

    all_metadata = []
    for item_id in item_ids:
        if require_attachment:
            c.execute("SELECT key FROM items WHERE itemID=?", (item_id,))
            row_key = c.fetchone()
            if not row_key or not row_key[0]:
                # skip if no file
                continue

        # 3) key -> file path or folder
        c.execute("SELECT key FROM items WHERE itemID=?", (item_id,))
        row_key = c.fetchone()
        doc_key = row_key[0] if row_key else ""

        # 4) Title (fieldID=110)
        c.execute("""
            SELECT itemDataValues.value
            FROM itemData
            JOIN itemDataValues ON itemData.valueID=itemDataValues.valueID
            WHERE itemData.itemID=? AND itemData.fieldID=110
        """, (item_id,))
        row_t = c.fetchone()
        title = row_t[0] if row_t else ""

        # 5) Authors
        c.execute("""
            SELECT group_concat(
                CASE WHEN creatorData.lastName IS NOT NULL AND creatorData.lastName!=''
                     THEN creatorData.lastName ELSE creatorData.lastName
                END, '; '
            ) AS authors
            FROM itemCreators
            JOIN creators     ON itemCreators.creatorID=creators.creatorID
            JOIN creatorData  ON creators.creatorDataID=creatorData.creatorDataID
            WHERE itemCreators.itemID=?
        """, (item_id,))
        row_a = c.fetchone()
        authors = row_a[0] if row_a else ""

        # 6) Year (fieldID=115)
        c.execute("""
            SELECT itemDataValues.value
            FROM itemData
            JOIN itemDataValues ON itemData.valueID=itemDataValues.valueID
            WHERE itemData.itemID=? AND itemData.fieldID=115
        """, (item_id,))
        row_y = c.fetchone()
        date_val = row_y[0] if row_y else ""
        year = ""
        if date_val:
            m = re.match(r"(\d{4})", date_val)
            if m: 
                year = m.group(1)

        meta_dict = {
            "itemID": item_id,
            "title": title,
            "authors": authors,
            "year": year,
            "key": doc_key
        }
        all_metadata.append(meta_dict)

    conn.close()
    return all_metadata


def sync_folder_with_db(folder_path, db_path):
    """
    Creates or updates the Zotero DB to represent:
      - One root-level collection for the main folder itself,
      - Additional collections for any subfolders,
      - Items for each file in each folder.

    This version ensures the *root* folder becomes a collection as well.
    """
    folder_path = Path(folder_path).resolve()
    if not folder_path.is_dir():
        raise NotADirectoryError(f"{folder_path} is not a valid directory.")

    conn = sqlite3.connect(str(db_path))
    c = conn.cursor()

    # -------------------------------------------------------
    # 0) Load existing collections/items
    # -------------------------------------------------------
    folder_to_coll = _load_existing_collections(c)
    file_to_item   = _load_existing_items(c)

    disk_folders   = set()
    disk_files     = set()

    for dirpath, dirnames, filenames in os.walk(folder_path):
      disk_folders.add(Path(dirpath).resolve())
      for f in filenames:
        if f.endswith(".npy"):
            continue  # skip .npy files
        if f.startswith(".") or f.startswith("~$"):
            continue  # skip hidden and Word temp lock files
        disk_files.add(Path(dirpath, f).resolve())


    # -------------------------------------------------------
    # (New) Always create main folder as a Collection
    # -------------------------------------------------------
    main_folder = folder_path.resolve()
    if main_folder not in folder_to_coll:
        # We have no existing Collection for the root folder
        # so create it with no parent (parent_coll_id=None).
        coll_id = _create_collection(conn, c,
                                     collection_name = main_folder.name,
                                     parent_coll_id  = None,
                                     full_path       = str(main_folder))
        folder_to_coll[main_folder] = coll_id

    # -------------------------------------------------------
    # 1) Detect removed folders
    # -------------------------------------------------------
    removed_folders = set(folder_to_coll.keys()) - disk_folders
    for old_folder in removed_folders:
        coll_id = folder_to_coll[old_folder]
        _remove_collection(conn, c, coll_id)
        del folder_to_coll[old_folder]

    # -------------------------------------------------------
    # 2) Detect new folders
    # -------------------------------------------------------
    new_folders = disk_folders - set(folder_to_coll.keys())
    new_folders_sorted = sorted(list(new_folders), key=lambda p: len(str(p)))
    for new_folder in new_folders_sorted:
        parent_folder = new_folder.parent
        if parent_folder in folder_to_coll:
            parent_coll_id = folder_to_coll[parent_folder]
        else:
            parent_coll_id = None
        coll_id = _create_collection(
            conn, c,
            collection_name = new_folder.name,
            parent_coll_id  = parent_coll_id,
            full_path       = str(new_folder)
        )
        folder_to_coll[new_folder] = coll_id

    # -------------------------------------------------------
    # 3) Detect removed files
    # -------------------------------------------------------
    removed_files = set(file_to_item.keys()) - disk_files
    for old_file in removed_files:
        item_id = file_to_item[old_file]
        _remove_item(conn, c, item_id)
        del file_to_item[old_file]

    # -------------------------------------------------------
    # 4) Detect new files
    # -------------------------------------------------------
    new_files = disk_files - set(file_to_item.keys())
    for new_file in new_files:
        parent_coll_id = None
        parent_folder  = new_file.parent
        if parent_folder in folder_to_coll:
            parent_coll_id = folder_to_coll[parent_folder]
        item_id = _create_item(conn, c, new_file, parent_coll_id)
        file_to_item[new_file] = item_id

    conn.commit()
    conn.close()
    print("Sync complete (root folder also stored as a top-level collection).")

def get_all_items(db_path):
    """
    Returns metadata for all items (itemTypeID != 14).
    Similar to above, but for entire 'items' table.
    """
    if not os.path.exists(db_path):
        raise FileNotFoundError(f"Zotero database not found: {db_path}")

    conn = sqlite3.connect(db_path)
    c = conn.cursor()
    c.execute("SELECT itemID FROM items WHERE itemTypeID!=14")
    rows = c.fetchall()
    item_ids = [r[0] for r in rows]

    all_metadata = []
    for item_id in item_ids:
        c.execute("SELECT key FROM items WHERE itemID=?", (item_id,))
        row_key = c.fetchone()
        doc_key = row_key[0] if row_key else ""

        # Title
        c.execute("""
            SELECT itemDataValues.value
            FROM itemData
            JOIN itemDataValues ON itemData.valueID=itemDataValues.valueID
            WHERE itemData.itemID=? AND itemData.fieldID=110
        """, (item_id,))
        row_t = c.fetchone()
        title = row_t[0] if row_t else ""

        # Authors
        c.execute("""
            SELECT group_concat(
                CASE WHEN creatorData.lastName IS NOT NULL AND creatorData.lastName!=''
                     THEN creatorData.lastName ELSE creatorData.lastName
                END, '; '
            )
            FROM itemCreators
            JOIN creators     ON itemCreators.creatorID=creators.creatorID
            JOIN creatorData  ON creators.creatorDataID=creatorData.creatorDataID
            WHERE itemCreators.itemID=?
        """, (item_id,))
        row_a = c.fetchone()
        authors = row_a[0] if row_a else ""

        # Year
        c.execute("""
            SELECT itemDataValues.value
            FROM itemData
            JOIN itemDataValues ON itemData.valueID=itemDataValues.valueID
            WHERE itemData.itemID=? AND itemData.fieldID=115
        """, (item_id,))
        row_y = c.fetchone()
        date_val = row_y[0] if row_y else ""
        year = ""
        if date_val:
            m = re.match(r"(\d{4})", date_val)
            if m: 
                year = m.group(1)

        meta_dict = {
            "itemID": item_id,
            "title": title,
            "authors": authors,
            "year": year,
            "key": doc_key
        }
        all_metadata.append(meta_dict)

    conn.close()
    return all_metadata

def _initialize_default_library(conn):
    c = conn.cursor()
    c.execute("SELECT libraryID FROM libraries WHERE libraryID=1")
    row = c.fetchone()
    if not row:
        c.execute("INSERT INTO libraries (libraryID, libraryType) VALUES (1,1)")
        conn.commit()
        


def initialize_zotero_db_from_skeleton(skeleton_db, target_db):
    """
    If target_db does not exist, copy skeleton_db to create it.
    Then ensure libraryID=1 exists in libraries table.
    Also ensure a default Global project exists as an entity.
    """
    from pathlib import Path
    import shutil
    import datetime

    target_db = Path(target_db)
    if not target_db.exists():
        shutil.copyfile(skeleton_db, target_db)
        print(f"Copied skeleton DB to {target_db}")

    conn = sqlite3.connect(str(target_db))
    try:
        _initialize_default_library(conn)

        # Ensure Global project exists as an entity
        c = conn.cursor()
        c.execute("SELECT entity_id FROM entities WHERE entity_name = 'Global' AND entity_type = 'project'")
        if not c.fetchone():
            # Step 1: Insert entity
            c.execute("""
                INSERT INTO entities (entity_name, entity_type)
                VALUES (?, 'project')
            """, ("Global",))
            global_id = c.lastrowid

            # Step 2: Tag the entity with metadata (created_by, created_at)
            #now = datetime.datetime.utcnow().isoformat()
            now = datetime.datetime.now(datetime.timezone.utc).isoformat()

            c.executemany("""
                INSERT INTO entity_tags (entity_id, tagCategory, tagValue)
                VALUES (?, ?, ?)
            """, [
                (global_id, "created_by", "system"),
                (global_id, "created_at", now),
                (global_id, "description", "System-level project for coordination")
            ])
            conn.commit()
            print("Inserted default Global project into entities + tags.")

    finally:
        conn.close()



from datetime import datetime
import sqlite3

def insert_project_note(db_path, project_name, note_text, created_by="system"):
    """
    Insert a note into the items table, tagged as a project note.
    If the project doesn't exist as an entity, create it.
    """
    conn = sqlite3.connect(db_path)
    c = conn.cursor()

    try:
        # Lookup or insert the project entity
        c.execute("SELECT entity_id FROM entities WHERE entity_name = ? AND entity_type = 'project'", (project_name,))
        row = c.fetchone()
        if row:
            project_id = row[0]
        else:
            c.execute("""
                INSERT INTO entities (entity_name, entity_type)
                VALUES (?, 'project')
            """, (project_name,))
            project_id = c.lastrowid

            now = datetime.utcnow().isoformat()
            c.executemany("""
                INSERT INTO entity_tags (entity_id, tagCategory, tagValue)
                VALUES (?, ?, ?)
            """, [
                (project_id, "created_by", created_by),
                (project_id, "created_at", now),
                (project_id, "description", f"Auto-created project: {project_name}")
            ])

        # Insert the note into items
        now = datetime.utcnow().isoformat()
        note_key = f"note_{project_name.lower().replace(' ', '_')}_{now}"
        c.execute("""
            INSERT INTO items (key, content, created_at)
            VALUES (?, ?, ?)
        """, (note_key, note_text, now))
        item_id = c.lastrowid

        # Tag the item to indicate project association and note type
        c.executemany("""
            INSERT INTO itemTags (itemID, tagCategory, tagResponse, tagSource)
            VALUES (?, ?, ?, ?)
        """, [
            (item_id, 'item_type', 'project_note', created_by),
            (item_id, 'project', project_name, created_by),
            (item_id, 'note_content', note_text, created_by),
            (item_id, 'created_by', created_by, created_by)
        ])

        conn.commit()
        print(f"Inserted note into project '{project_name}' as item {item_id}.")

    finally:
        conn.close()



def _load_existing_collections(c):
    folder_map = {}
    c.execute("SELECT collectionID, collectionName, key FROM collections")
    for row in c.fetchall():
        coll_id, coll_name, folder_path_str = row
        if folder_path_str and os.path.exists(folder_path_str):
            folder_map[Path(folder_path_str).resolve()] = coll_id
    return folder_map

def _load_existing_items(c):
    file_map = {}
    c.execute("SELECT itemID, key FROM items WHERE itemTypeID!=14")
    for row in c.fetchall():
        item_id, file_path_str = row
        if file_path_str and os.path.exists(file_path_str):
            file_map[Path(file_path_str).resolve()] = item_id
    return file_map

def _remove_collection(conn, c, coll_id):
    c.execute("DELETE FROM collectionItems WHERE collectionID=?", (coll_id,))
    c.execute("DELETE FROM collections WHERE collectionID=?", (coll_id,))
    conn.commit()

def _remove_item(conn, c, item_id):
    c.execute("DELETE FROM itemData WHERE itemID=?", (item_id,))
    c.execute("DELETE FROM itemCreators WHERE itemID=?", (item_id,))
    c.execute("DELETE FROM collectionItems WHERE itemID=?", (item_id,))
    c.execute("DELETE FROM items WHERE itemID=?", (item_id,))
    conn.commit()

def _create_collection(conn, c, collection_name, parent_coll_id, full_path):
    library_id = 1
    c.execute("""
        INSERT INTO collections (libraryID, collectionName, parentCollectionID, key)
        VALUES (?,?,?,?)
    """, (library_id, collection_name, parent_coll_id, full_path))
    conn.commit()
    return c.lastrowid

def _create_item(conn, c, file_path, coll_id):
    random_key = _generate_random_key()
    library_id = 1
    c.execute("""
        INSERT INTO items (libraryID, itemTypeID, key)
        VALUES (?,?,?)
    """, (library_id, 2, str(file_path)))
    item_id = c.lastrowid

    if coll_id is not None:
        c.execute("INSERT INTO collectionItems (collectionID, itemID) VALUES (?,?)", (coll_id, item_id))

    filename = Path(file_path).name
    c.execute("SELECT valueID FROM itemDataValues WHERE value=?", (filename,))
    row_val = c.fetchone()
    if row_val:
        value_id_title = row_val[0]
    else:
        c.execute("INSERT INTO itemDataValues (value) VALUES (?)",(filename,))
        value_id_title = c.lastrowid

    # fieldID=110 => Title
    c.execute("INSERT INTO itemData (itemID,fieldID,valueID) VALUES (?,?,?)",(item_id,110,value_id_title))

    # maybe parse year
    year_match = re.search(r"\b(19\d{2}|20\d{2})\b", filename)
    if year_match:
        year_str = year_match.group(1)
        c.execute("SELECT valueID FROM itemDataValues WHERE value=?",(year_str,))
        row_yr = c.fetchone()
        if row_yr:
            value_id_year = row_yr[0]
        else:
            c.execute("INSERT INTO itemDataValues (value) VALUES (?)",(year_str,))
            value_id_year = c.lastrowid

        # fieldID=115 => date/year
        c.execute("INSERT INTO itemData (itemID,fieldID,valueID) VALUES (?,?,?)",(item_id,115,value_id_year))

    conn.commit()
    return item_id

def _generate_random_key():
    import uuid
    return uuid.uuid4().hex[:8]

###############################################################################
# 2) snippet-based text searching
###############################################################################

def insert_search_results(db_path, results, queryID, query_text, collection_name):
    """
    Insert a batch of vector search results into the search_results table.
    
    Each `res` dict must contain: snippetID, matched_word, context
    """
    conn = sqlite3.connect(db_path)
    c = conn.cursor()
    for res in results:
        c.execute("""
            INSERT INTO search_results (queryID, snippetID, query, matched_word, context, collection_name)
            VALUES (?, ?, ?, ?, ?, ?)
        """, (queryID, res["snippetID"], query_text, res["matched_word"], res["context"], collection_name))
    conn.commit()
    conn.close()

def get_search_results(db_path, collection_name=None):
    """
    Retrieve search results with document name, based on snippetID → itemID → items.key.
    """
    conn = sqlite3.connect(db_path)
    c = conn.cursor()

    if collection_name:
        c.execute("""
            SELECT sr.snippetID, sr.query, sr.matched_word, sr.context, sr.collection_name, sr.timestamp,
                   i.key AS document
            FROM search_results sr
            JOIN documentEmbeddings de ON sr.snippetID = de.snippetID
            JOIN items i ON de.itemID = i.itemID
            WHERE sr.collection_name = ?
            ORDER BY sr.timestamp DESC
        """, (collection_name,))
    else:
        c.execute("""
            SELECT sr.snippetID, sr.query, sr.matched_word, sr.context, sr.collection_name, sr.timestamp,
                   i.key AS document
            FROM search_results sr
            JOIN documentEmbeddings de ON sr.snippetID = de.snippetID
            JOIN items i ON de.itemID = i.itemID
            ORDER BY sr.timestamp DESC
        """)

    rows = c.fetchall()
    conn.close()

    df = pd.DataFrame(rows, columns=[
        "snippetID", "query", "matched_word", "context", "collection_name", "timestamp", "document"
    ])
    df["document"] = df["document"].apply(lambda x: os.path.basename(x) if x else "(No Name)")
    return df


def get_entity_id(db_path, name, entity_type):
    with sqlite3.connect(db_path) as con:
        cur = con.cursor()
        cur.execute("SELECT entity_id FROM entities WHERE entity_name = ? AND entity_type = ?", (name, entity_type))
        result = cur.fetchone()
        return result[0] if result else None

def get_entity_tags(db_path, name, entity_type):
    with sqlite3.connect(db_path) as con:
        cur = con.cursor()
        cur.execute("""
            SELECT tagCategory, tagValue 
            FROM entity_tags 
            WHERE entity_id = (SELECT entity_id FROM entities WHERE entity_name = ? AND entity_type = ?)
        """, (name, entity_type))
        return cur.fetchall()

def add_entity_tag(db_path, name, entity_type, tagCategory, tagValue):
    entity_id = get_entity_id(db_path, name, entity_type)
    if entity_id is not None:
        with sqlite3.connect(db_path) as con:
            cur = con.cursor()
            cur.execute("""
                INSERT INTO entity_tags (entity_id, tagCategory, tagValue)
                VALUES (?, ?, ?)""", (entity_id, tagCategory, tagValue))



###############################################################################
# 3) DocumentEmbeddings table -> no chunkText
###############################################################################

def split_text_into_chunks(text, chunk_size=25):
    """
    Returns a list of (start_word_idx, end_word_idx) for each chunk.

    e.g. if text has 103 words and chunk_size=25,
         you'll get 5 chunks: 
           chunk0 => words [0..24]
           chunk1 => words [25..49]
           chunk2 => ...
    """
    import re
    words = re.split(r"\s+", text.strip())
    chunks = []
    start_idx = 0
    while start_idx < len(words):
        end_idx = start_idx + chunk_size
        if end_idx > len(words):
            end_idx = len(words)
        chunks.append((start_idx, end_idx))
        start_idx += chunk_size
    return words, chunks


def generate_document_embeddings(db_path, chunk_size=25, model_name="sentence-transformers/all-MiniLM-L6-v2"):
    """
    1) Ensures 'documentEmbeddings' table
    2) For each item in 'items', read text, chunk by chunk_size words
    3) For each chunk, store snippetID row with (chunkStart, chunkEnd, embeddingModel, chunkSize)
    4) Also save a .npy file with the embedding in e.g. [db_folder]/embedding_data/snippet_{snippetID}.npy
    """
    import os
    from pathlib import Path
    from sentence_transformers import SentenceTransformer

    # -------------------------------------------------------
    # Step 1: Ensure caching folder + load model from cache
    # -------------------------------------------------------
    CACHE_DIR = str(Path.home() / ".cache" / "sentence_transformers")

    try:
        embedder = SentenceTransformer(model_name, cache_folder=CACHE_DIR)
    except Exception as e:
        raise RuntimeError(f"Failed to load embedding model '{model_name}': {e}")

    # -------------------------------------------------------
    # Step 2: Prepare embedding table and paths
    # -------------------------------------------------------

    db_folder = os.path.dirname(db_path)
    emb_folder = os.path.join(db_folder, "embedding_data")
    os.makedirs(emb_folder, exist_ok=True)

    conn = sqlite3.connect(db_path)
    c = conn.cursor()

    # -------------------------------------------------------
    # Step 3: Get all items in DB and generate embeddings
    # -------------------------------------------------------
    c.execute("SELECT itemID, key FROM items WHERE itemTypeID!=14")
    rows = c.fetchall()

    for (item_id, file_key) in rows:
        if file_key and os.path.exists(file_key):
            text = read_text_file(file_key)
            if not text or not text.strip():
                continue

            words, chunk_ranges = split_text_into_chunks(text, chunk_size)
            chunk_texts = [" ".join(words[start_i:end_i]) for (start_i, end_i) in chunk_ranges]

            # Embed chunks
            embeddings = embedder.encode(chunk_texts, convert_to_numpy=True)

            # Store embeddings and metadata
            for idx, (start_i, end_i) in enumerate(chunk_ranges):
                c.execute("""
                  INSERT INTO documentEmbeddings
                    (itemID, chunkIndex, chunkStart, chunkEnd, embeddingModel, chunkSize)
                  VALUES (?,?,?,?,?,?)
                """, (item_id, idx, start_i, end_i, model_name, chunk_size))
                snippet_id = c.lastrowid

                emb_path = os.path.join(emb_folder, f"snippet_{snippet_id}.npy")
                np.save(emb_path, embeddings[idx])

    conn.commit()
    conn.close()
    print("Done generating document embeddings in table + .npy files.")



###############################################################################
# 4) Document Handeling
###############################################################################

from docx import Document
from pptx import Presentation

def extract_text_from_docx(path):
    doc = Document(path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(path):
    prs = Presentation(path)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

